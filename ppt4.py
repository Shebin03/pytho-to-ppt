from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
import json

prs = Presentation()

title_only_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_only_slide_layout)
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
shapes = slide.shapes
#shapes.title.text = 'Profile Automation'
cols = 2
rows=7
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(5.0)
table = shapes.add_table(rows, cols, left, top, width, height).table
# set column widths

table.columns[0].width = Inches(5.0)
table.columns[1].width = Inches(5.0)
table.columns[1].height = Inches(5.0)
table.rows[1].height = Inches(0.4)
table.rows[2].height = Inches(1.5)
table.rows[3].height = Inches(0.4)
table.rows[4].height = Inches(1.5)
table.rows[5].height = Inches(0.4)
table.rows[6].height = Inches(1.5)

# importing json
data2 = json.load(open('ppt1.json'))
# data1=json.dumps(data)
# data2=json.loads(data1)


# write column headings
table.cell(0, 0).text = 'Shebin Paul\nSoftware Engineering'

# write body cells
table.cell(1, 0).text = 'Summary'
table.cell(1, 1).text = 'Key Projects'
table.cell(3, 0).text = 'Academic Details'
table.cell(5, 0).text = 'Area of expertise'

table.cell(2, 0).text = data2["Summary"]
table.cell(4, 0).text = data2["Academic Details"]
table.cell(5, 0).text = data2["Area of expertise"]
table.cell(2, 1).text = data2["Key Projects"]

# applying the colors to cell

cell = table.cell(1, 0)  # where to apply.
fill = cell.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x41, 0x69, 0xe1)

cell = table.cell(3, 0)
fill = cell.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x00)

cell = table.cell(5, 0)
fill = cell.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xb0, 0xe0, 0xe6)



cell = table.cell(1, 1)
fill = cell.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xff, 0x7f, 0x00)

# merging the cells

cell = table.cell(0, 0)
other_cell = table.cell(0, 1)
cell.merge(other_cell)

cell = table.cell(2, 1)
other_cell = table.cell(6, 1)
cell.merge(other_cell)


prs.save('test.pptx')
